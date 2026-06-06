#!/usr/bin/env python3
"""Create module-level cutouts from image-only PPT slides.

Subcommands:
  prepare: extract full-slide images and auto-detect card/module boxes
  serve:   launch a browser-adjustable red-box review page
  review:  prepare, then serve
  build:   create a PPTX with a white-filled base layer plus module image layers
"""

from __future__ import annotations

import argparse
import json
import mimetypes
import re
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from io import BytesIO
from pathlib import Path
from typing import Any, Iterable

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Emu

try:
    import cv2  # type: ignore
    import numpy as np  # type: ignore
except Exception:  # pragma: no cover - optional fallback
    cv2 = None
    np = None


REVIEW_HTML = r"""<!doctype html>
<html lang="zh-CN">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<link rel="icon" href="data:,">
<title>Module Cutout Review</title>
<style>
  :root { --ink:#0b2146; --muted:#5f6d82; --line:#d7dee9; --panel:#f6f8fb; --red:#d62d20; --blue:#0c4a86; }
  * { box-sizing:border-box; }
  body { margin:0; font-family:"Microsoft YaHei","Segoe UI",Arial,sans-serif; color:var(--ink); background:#eef2f7; }
  .topbar { position:sticky; top:0; z-index:20; display:flex; align-items:center; justify-content:space-between; gap:12px; padding:10px 14px; background:#fff; border-bottom:1px solid var(--line); }
  .title { min-width:220px; font-weight:700; font-size:16px; white-space:nowrap; }
  .actions { display:flex; gap:8px; flex-wrap:wrap; justify-content:flex-end; }
  button, select { border:1px solid #b7c3d4; background:#fff; color:var(--ink); border-radius:6px; padding:7px 10px; font:inherit; cursor:pointer; }
  button.primary { background:var(--blue); border-color:var(--blue); color:#fff; }
  button.danger { border-color:#e0a19b; color:#9c1d14; }
  button:disabled { cursor:not-allowed; opacity:.45; }
  .status { min-width:150px; color:var(--muted); font-size:13px; text-align:right; }
  .layout { display:grid; grid-template-columns:250px minmax(0,1fr) 250px; gap:12px; padding:12px; min-height:calc(100vh - 56px); }
  .side, .inspector { background:#fff; border:1px solid var(--line); border-radius:8px; overflow:hidden; min-height:0; }
  .side h2, .inspector h2 { margin:0; padding:10px 12px; font-size:14px; border-bottom:1px solid var(--line); background:var(--panel); }
  .slide-list { height:calc(100vh - 112px); overflow:auto; padding:8px; display:grid; gap:8px; }
  .slide-card { border:1px solid var(--line); border-radius:7px; padding:6px; background:#fff; cursor:pointer; }
  .slide-card.active { border-color:var(--blue); box-shadow:0 0 0 2px rgba(12,74,134,.15); }
  .slide-card img { width:100%; display:block; border-radius:4px; background:#eef2f7; }
  .slide-meta { display:flex; justify-content:space-between; margin-top:5px; font-size:12px; color:var(--muted); }
  .stage-wrap { min-width:0; display:flex; align-items:flex-start; justify-content:center; overflow:auto; background:#dfe6ef; border:1px solid var(--line); border-radius:8px; padding:14px; height:calc(100vh - 80px); }
  .stage { position:relative; background:#fff; box-shadow:0 12px 34px rgba(6,22,48,.15); line-height:0; user-select:none; }
  .stage img { display:block; width:min(100%, calc((100vh - 120px) * 16 / 9)); max-width:100%; height:auto; }
  .box { position:absolute; border:2px solid var(--red); background:rgba(214,45,32,.04); cursor:move; min-width:12px; min-height:12px; }
  .box.selected { border-width:3px; background:rgba(214,45,32,.1); box-shadow:0 0 0 2px rgba(255,255,255,.8); }
  .box .label { position:absolute; top:-22px; left:-2px; background:var(--red); color:#fff; font-size:12px; line-height:18px; height:18px; padding:0 6px; border-radius:4px 4px 0 0; }
  .handle { position:absolute; width:12px; height:12px; background:#fff; border:2px solid var(--red); border-radius:50%; }
  .handle.nw { left:-7px; top:-7px; cursor:nwse-resize; }
  .handle.ne { right:-7px; top:-7px; cursor:nesw-resize; }
  .handle.sw { left:-7px; bottom:-7px; cursor:nesw-resize; }
  .handle.se { right:-7px; bottom:-7px; cursor:nwse-resize; }
  .inspector-body { height:calc(100vh - 112px); overflow:auto; padding:10px; display:grid; gap:10px; align-content:start; }
  .kv { display:grid; grid-template-columns:1fr 1fr; gap:8px; }
  label { display:grid; gap:4px; font-size:12px; color:var(--muted); }
  input { width:100%; border:1px solid #b7c3d4; border-radius:6px; padding:7px 8px; font:inherit; color:var(--ink); }
  .box-list { display:grid; gap:6px; }
  .box-row { display:flex; justify-content:space-between; gap:8px; align-items:center; border:1px solid var(--line); border-radius:6px; padding:7px 8px; background:#fff; cursor:pointer; font-size:13px; }
  .box-row.active { border-color:var(--red); background:#fff6f5; }
  .hint { color:var(--muted); font-size:12px; line-height:1.55; margin:0; }
  @media (max-width:1100px) { .layout { grid-template-columns:180px minmax(0,1fr); } .inspector { grid-column:1 / -1; } .inspector-body { height:auto; } }
</style>
</head>
<body>
<div class="topbar">
  <div class="title">模块红框校准</div>
  <div class="actions">
    <button id="prevBtn">上一页</button><select id="slideSelect"></select><button id="nextBtn">下一页</button>
    <button id="addBtn">新增框</button><button id="deleteBtn" class="danger">删除框</button>
    <button id="saveBtn" class="primary">保存调整</button><button id="downloadBtn">下载 JSON</button>
  </div>
  <div id="status" class="status">加载中</div>
</div>
<div class="layout">
  <aside class="side"><h2>幻灯片</h2><div id="slideList" class="slide-list"></div></aside>
  <main class="stage-wrap"><div id="stage" class="stage"><img id="slideImage" alt="slide"></div></main>
  <aside class="inspector"><h2>当前红框</h2><div class="inspector-body">
    <p class="hint">拖动红框移动，拖四角圆点缩放。方向键微调，Shift + 方向键每次移动 10 像素，Ctrl+S 保存。</p>
    <div class="kv"><label>X<input id="xInput" type="number"></label><label>Y<input id="yInput" type="number"></label><label>W<input id="wInput" type="number"></label><label>H<input id="hInput" type="number"></label></div>
    <div id="boxList" class="box-list"></div>
  </div></aside>
</div>
<script>
const state = { slides: [], current: 0, selected: 0, scale: 1, dirty: false };
window.reviewState = state;
const el = (id) => document.getElementById(id);
const stage = el("stage");
const slideImage = el("slideImage");
const status = el("status");
function setStatus(text) { status.textContent = text; }
function slide() { return state.slides[state.current]; }
function selectedBox() { const s = slide(); return s && s.boxes[state.selected] ? s.boxes[state.selected] : null; }
function markDirty() { state.dirty = true; setStatus("未保存"); renderSidebar(); }
function clampBox(box, s) {
  box.w = Math.max(8, Math.round(box.w)); box.h = Math.max(8, Math.round(box.h));
  box.x = Math.max(0, Math.min(Math.round(box.x), s.size[0] - box.w));
  box.y = Math.max(0, Math.min(Math.round(box.y), s.size[1] - box.h));
}
function toScreen(box) { return { left: box.x * state.scale, top: box.y * state.scale, width: box.w * state.scale, height: box.h * state.scale }; }
function fromScreenDelta(dx, dy) { return { dx: dx / state.scale, dy: dy / state.scale }; }
function render() {
  const s = slide(); if (!s) return;
  slideImage.src = "/slides/slide_" + String(s.slide).padStart(2, "0") + ".png?v=" + Date.now();
  slideImage.onload = () => { state.scale = slideImage.clientWidth / s.size[0]; drawBoxes(); };
  el("slideSelect").value = String(state.current);
  renderSidebar(); renderInspector(); setStatus(state.dirty ? "未保存" : "已加载");
}
function drawBoxes() {
  stage.querySelectorAll(".box").forEach((node) => node.remove());
  const s = slide(); state.scale = slideImage.clientWidth / s.size[0];
  s.boxes.forEach((box, index) => {
    const rect = toScreen(box);
    const node = document.createElement("div");
    node.className = "box" + (index === state.selected ? " selected" : "");
    node.style.left = rect.left + "px"; node.style.top = rect.top + "px"; node.style.width = rect.width + "px"; node.style.height = rect.height + "px";
    node.dataset.index = String(index);
    node.innerHTML = '<div class="label">' + (index + 1) + '</div><div class="handle nw"></div><div class="handle ne"></div><div class="handle sw"></div><div class="handle se"></div>';
    node.addEventListener("pointerdown", startPointer);
    stage.appendChild(node);
  });
}
function startPointer(event) {
  event.preventDefault();
  const node = event.currentTarget; const index = Number(node.dataset.index); state.selected = index;
  const s = slide(); const box = s.boxes[index];
  const handle = event.target.classList.contains("handle") ? [...event.target.classList].find((c) => ["nw","ne","sw","se"].includes(c)) : null;
  const start = { x: event.clientX, y: event.clientY, box: { ...box } };
  node.setPointerCapture(event.pointerId);
  function move(e) {
    const d = fromScreenDelta(e.clientX - start.x, e.clientY - start.y);
    if (!handle) { box.x = start.box.x + d.dx; box.y = start.box.y + d.dy; }
    else {
      let x1 = start.box.x, y1 = start.box.y, x2 = start.box.x + start.box.w, y2 = start.box.y + start.box.h;
      if (handle.includes("w")) x1 += d.dx; if (handle.includes("e")) x2 += d.dx; if (handle.includes("n")) y1 += d.dy; if (handle.includes("s")) y2 += d.dy;
      if (x2 - x1 < 8) x2 = x1 + 8; if (y2 - y1 < 8) y2 = y1 + 8;
      box.x = x1; box.y = y1; box.w = x2 - x1; box.h = y2 - y1;
    }
    clampBox(box, s); drawBoxes(); renderInspector(); state.dirty = true; setStatus("未保存");
  }
  function up(e) { node.releasePointerCapture(e.pointerId); window.removeEventListener("pointermove", move); window.removeEventListener("pointerup", up); markDirty(); }
  window.addEventListener("pointermove", move); window.addEventListener("pointerup", up); drawBoxes(); renderInspector();
}
function renderSidebar() {
  const list = el("slideList"); list.innerHTML = "";
  state.slides.forEach((s, i) => {
    const card = document.createElement("div");
    card.className = "slide-card" + (i === state.current ? " active" : "");
    card.innerHTML = '<img src="/slides/slide_' + String(s.slide).padStart(2, "0") + '.png"><div class="slide-meta"><span>S' + String(s.slide).padStart(2, "0") + '</span><span>' + s.boxes.length + ' 个框</span></div>';
    card.addEventListener("click", () => { state.current = i; state.selected = 0; render(); });
    list.appendChild(card);
  });
}
function renderInspector() {
  const box = selectedBox(); const inputs = ["xInput","yInput","wInput","hInput"].map(el);
  inputs.forEach((input) => input.disabled = !box); el("deleteBtn").disabled = !box;
  if (box) { el("xInput").value = Math.round(box.x); el("yInput").value = Math.round(box.y); el("wInput").value = Math.round(box.w); el("hInput").value = Math.round(box.h); }
  else { inputs.forEach((input) => input.value = ""); }
  const list = el("boxList"); list.innerHTML = "";
  slide().boxes.forEach((b, i) => {
    const row = document.createElement("div");
    row.className = "box-row" + (i === state.selected ? " active" : "");
    row.innerHTML = '<strong>框 ' + (i + 1) + '</strong><span>' + Math.round(b.x) + ', ' + Math.round(b.y) + ', ' + Math.round(b.w) + ' x ' + Math.round(b.h) + '</span>';
    row.addEventListener("click", () => { state.selected = i; drawBoxes(); renderInspector(); });
    list.appendChild(row);
  });
}
function updateFromInputs() {
  const b = selectedBox(); if (!b) return;
  b.x = Number(el("xInput").value); b.y = Number(el("yInput").value); b.w = Number(el("wInput").value); b.h = Number(el("hInput").value);
  clampBox(b, slide()); drawBoxes(); renderInspector(); markDirty();
}
async function save() {
  const res = await fetch("/api/save", { method: "POST", headers: { "Content-Type": "application/json" }, body: JSON.stringify({ slides: state.slides }) });
  if (!res.ok) { setStatus("保存失败"); alert(await res.text()); return; }
  state.dirty = false; setStatus("已保存 adjusted_modules.json");
}
function addBox() {
  const s = slide(); const w = Math.round(s.size[0] * 0.12); const h = Math.round(s.size[1] * 0.16);
  s.boxes.push({ x: Math.round((s.size[0] - w) / 2), y: Math.round((s.size[1] - h) / 2), w, h });
  state.selected = s.boxes.length - 1; drawBoxes(); renderInspector(); markDirty();
}
function deleteBox() {
  const s = slide(); if (!selectedBox()) return;
  s.boxes.splice(state.selected, 1); state.selected = Math.max(0, Math.min(state.selected, s.boxes.length - 1)); drawBoxes(); renderInspector(); markDirty();
}
function downloadJson() {
  const blob = new Blob([JSON.stringify({ slides: state.slides }, null, 2)], { type: "application/json" });
  const url = URL.createObjectURL(blob); const a = document.createElement("a"); a.href = url; a.download = "adjusted_modules.json"; a.click(); URL.revokeObjectURL(url);
}
function moveSelected(dx, dy) {
  const b = selectedBox(); if (!b) return;
  b.x += dx; b.y += dy; clampBox(b, slide()); drawBoxes(); renderInspector(); markDirty();
}
function initControls() {
  el("prevBtn").onclick = () => { state.current = Math.max(0, state.current - 1); state.selected = 0; render(); };
  el("nextBtn").onclick = () => { state.current = Math.min(state.slides.length - 1, state.current + 1); state.selected = 0; render(); };
  el("addBtn").onclick = addBox; el("deleteBtn").onclick = deleteBox; el("saveBtn").onclick = save; el("downloadBtn").onclick = downloadJson;
  ["xInput","yInput","wInput","hInput"].forEach((id) => el(id).addEventListener("change", updateFromInputs));
  el("slideSelect").addEventListener("change", (event) => { state.current = Number(event.target.value); state.selected = 0; render(); });
  window.addEventListener("resize", drawBoxes);
  window.addEventListener("keydown", (event) => {
    if ((event.ctrlKey || event.metaKey) && event.key.toLowerCase() === "s") { event.preventDefault(); save(); return; }
    const step = event.shiftKey ? 10 : 1;
    if (event.key === "ArrowLeft") moveSelected(-step, 0);
    if (event.key === "ArrowRight") moveSelected(step, 0);
    if (event.key === "ArrowUp") moveSelected(0, -step);
    if (event.key === "ArrowDown") moveSelected(0, step);
  });
}
async function boot() {
  const res = await fetch("/api/manifest"); const payload = await res.json();
  state.slides = payload.slides.map((s) => ({ slide: s.slide, size: s.size, boxes: (s.boxes || []).map((b) => ({ x:b.x, y:b.y, w:b.w, h:b.h })) }));
  el("slideSelect").innerHTML = state.slides.map((s, i) => '<option value="' + i + '">第 ' + s.slide + ' 页 (' + s.boxes.length + ')</option>').join("");
  initControls(); state.current = 0; state.selected = 0; state.dirty = false; render(); window.setTimeout(() => { if (!state.dirty) setStatus("已加载"); }, 0);
}
boot().catch((error) => { console.error(error); setStatus("加载失败"); alert(error.message); });
</script>
</body>
</html>
"""


def read_json(path: Path) -> Any:
    return json.loads(path.read_text(encoding="utf-8"))


def write_json(path: Path, value: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(value, ensure_ascii=False, indent=2), encoding="utf-8")


def full_slide_picture(slide: Any) -> Any | None:
    pictures = []
    for shape in slide.shapes:
        if hasattr(shape, "image"):
            try:
                pictures.append((int(shape.width) * int(shape.height), shape))
            except Exception:
                pass
    if not pictures:
        return None
    return max(pictures, key=lambda item: item[0])[1]


def edge_score(mask: Any, x: int, y: int, w: int, h: int) -> float:
    crop = mask[y : y + h, x : x + w]
    if crop.size == 0:
        return 0.0
    t = max(3, min(22, h // 16, w // 16))
    values = (crop[:t, :].mean(), crop[-t:, :].mean(), crop[:, :t].mean(), crop[:, -t:].mean())
    return float(min(values) / 255.0)


def detect_module_boxes(image: Image.Image) -> list[dict[str, int | float]]:
    if cv2 is None or np is None:
        return []
    arr = np.array(image.convert("RGB"))
    height, width = arr.shape[:2]
    r, g, b = arr[:, :, 0], arr[:, :, 1], arr[:, :, 2]
    dark = (r < 95) & (g < 125) & (b < 175) & ((b > r + 10) | (r < 55))
    ymask = np.zeros((height, width), dtype=bool)
    ymask[int(height * 0.20) : int(height * 0.88), :] = True
    mask = (dark & ymask).astype("uint8") * 255
    kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (5, 5))
    mask2 = cv2.morphologyEx(mask, cv2.MORPH_CLOSE, kernel, iterations=1)
    contours, _ = cv2.findContours(mask2, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    boxes = []
    for contour in contours:
        x, y, w, h = cv2.boundingRect(contour)
        area = w * h
        if area < width * height * 0.002 or area > width * height * 0.20:
            continue
        if w < width * 0.035 or h < height * 0.07:
            continue
        if w > width * 0.55 or h > height * 0.62:
            continue
        aspect = w / h
        if aspect < 0.18 or aspect > 4.5:
            continue
        score = edge_score(mask2, x, y, w, h)
        if score < 0.006:
            continue
        crop = arr[y : y + h, x : x + w]
        light = ((crop[:, :, 0] > 205) & (crop[:, :, 1] > 205) & (crop[:, :, 2] > 205)).mean()
        if light < 0.25:
            continue
        boxes.append([x, y, w, h, score, float(light)])

    boxes = sorted(boxes, key=lambda item: item[2] * item[3], reverse=True)
    kept = []
    for box in boxes:
        x, y, w, h = box[:4]
        contained = False
        for prior in kept:
            x2, y2, w2, h2 = prior[:4]
            ix = max(0, min(x + w, x2 + w2) - max(x, x2))
            iy = max(0, min(y + h, y2 + h2) - max(y, y2))
            if (ix * iy) / (w * h) > 0.72:
                contained = True
                break
        if not contained:
            kept.append(box)

    pad = max(10, int(min(width, height) * 0.008))
    result = []
    for x, y, w, h, score, light in sorted(kept, key=lambda item: (item[1], item[0])):
        x0 = max(0, x - pad)
        y0 = max(0, y - pad)
        x1 = min(width, x + w + pad)
        y1 = min(height, y + h + pad)
        result.append({"x": x0, "y": y0, "w": x1 - x0, "h": y1 - y0, "score": score, "light": light})
    return result


def natural_sort_key(path: Path) -> list[Any]:
    return [int(part) if part.isdigit() else part.lower() for part in re.split(r"(\d+)", path.name)]


def prepare_review(pptx: Path, work_dir: Path, force: bool = False) -> Path:
    slides_dir = work_dir / "slides"
    manifest_path = work_dir / "manifest.json"
    if manifest_path.exists() and not force:
        return manifest_path

    slides_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(pptx))
    manifest = []
    for slide_no, slide in enumerate(prs.slides, 1):
        picture = full_slide_picture(slide)
        if picture is None:
            raise RuntimeError(f"slide {slide_no} has no picture source")
        image = Image.open(BytesIO(picture.image.blob)).convert("RGB")
        slide_path = slides_dir / f"slide_{slide_no:02d}.png"
        image.save(slide_path)
        manifest.append({"slide": slide_no, "size": list(image.size), "boxes": detect_module_boxes(image)})
    write_json(manifest_path, {"source": str(pptx), "slides": manifest})
    return manifest_path


def normalize_payload(payload: Any, manifest_path: Path) -> dict[str, Any]:
    if not isinstance(payload, dict) or not isinstance(payload.get("slides"), list):
        raise ValueError("payload must contain a slides array")
    normalized = {"sourceManifest": str(manifest_path), "slides": []}
    for slide in payload["slides"]:
        slide_no = int(slide["slide"])
        size = [int(slide["size"][0]), int(slide["size"][1])]
        boxes = []
        for box in slide.get("boxes", []):
            x, y, w, h = round(float(box["x"])), round(float(box["y"])), round(float(box["w"])), round(float(box["h"]))
            if w <= 1 or h <= 1:
                raise ValueError(f"invalid box on slide {slide_no}")
            boxes.append({"x": x, "y": y, "w": w, "h": h})
        normalized["slides"].append({"slide": slide_no, "size": size, "boxes": boxes})
    return normalized


def serve_review(work_dir: Path, port: int) -> None:
    manifest_path = work_dir / "manifest.json"
    adjusted_path = work_dir / "adjusted_modules.json"
    slides_dir = work_dir / "slides"
    if not manifest_path.exists():
        raise FileNotFoundError(manifest_path)

    class Handler(BaseHTTPRequestHandler):
        def log_message(self, fmt: str, *args: Any) -> None:
            print(f"{self.address_string()} - {fmt % args}")

        def send_bytes(self, status: int, body: bytes, content_type: str) -> None:
            self.send_response(status)
            self.send_header("Content-Type", content_type)
            self.send_header("Cache-Control", "no-store")
            self.end_headers()
            self.wfile.write(body)

        def do_GET(self) -> None:  # noqa: N802
            if self.path == "/" or self.path.startswith("/?"):
                self.send_bytes(200, REVIEW_HTML.encode("utf-8"), "text/html; charset=utf-8")
                return
            if self.path.startswith("/api/manifest"):
                source = adjusted_path if adjusted_path.exists() else manifest_path
                payload = read_json(source)
                self.send_bytes(200, json.dumps(payload, ensure_ascii=False).encode("utf-8"), "application/json; charset=utf-8")
                return
            if self.path.startswith("/slides/"):
                name = Path(self.path.split("?", 1)[0]).name
                if not re.fullmatch(r"slide_\d{2}\.png", name):
                    self.send_bytes(404, b"not found", "text/plain; charset=utf-8")
                    return
                file = slides_dir / name
                if not file.exists():
                    self.send_bytes(404, b"not found", "text/plain; charset=utf-8")
                    return
                self.send_bytes(200, file.read_bytes(), mimetypes.guess_type(file.name)[0] or "image/png")
                return
            if self.path == "/favicon.ico":
                self.send_response(204)
                self.end_headers()
                return
            self.send_bytes(404, b"not found", "text/plain; charset=utf-8")

        def do_POST(self) -> None:  # noqa: N802
            if self.path != "/api/save":
                self.send_bytes(404, b"not found", "text/plain; charset=utf-8")
                return
            try:
                length = int(self.headers.get("Content-Length", "0"))
                payload = json.loads(self.rfile.read(length).decode("utf-8"))
                normalized = normalize_payload(payload, manifest_path)
                write_json(adjusted_path, normalized)
                self.send_bytes(200, json.dumps({"ok": True, "file": str(adjusted_path)}, ensure_ascii=False).encode("utf-8"), "application/json; charset=utf-8")
            except Exception as exc:
                self.send_bytes(400, str(exc).encode("utf-8"), "text/plain; charset=utf-8")

    server = ThreadingHTTPServer(("127.0.0.1", port), Handler)
    print(f"Module review: http://127.0.0.1:{port}/")
    print(f"Saving adjustments to: {adjusted_path}")
    print("Press Ctrl+C to stop.")
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        pass
    finally:
        server.server_close()


def clamp_box(box: dict[str, Any], width: int, height: int) -> tuple[int, int, int, int]:
    x = max(0, min(round(float(box["x"])), width - 1))
    y = max(0, min(round(float(box["y"])), height - 1))
    w = max(1, round(float(box["w"])))
    h = max(1, round(float(box["h"])))
    if x + w > width:
        w = width - x
    if y + h > height:
        h = height - y
    return x, y, w, h


def unique_output_path(path: Path, overwrite: bool) -> Path:
    if overwrite or not path.exists():
        return path
    stem, suffix = path.stem, path.suffix
    index = 2
    while True:
        candidate = path.with_name(f"{stem}_v{index}{suffix}")
        if not candidate.exists():
            return candidate
        index += 1


def add_picture(slide: Any, image_path: Path, x: int, y: int, w: int, h: int, img_w: int, img_h: int, prs: Presentation, name: str) -> None:
    pic = slide.shapes.add_picture(
        str(image_path),
        Emu(round(x / img_w * prs.slide_width)),
        Emu(round(y / img_h * prs.slide_height)),
        width=Emu(round(w / img_w * prs.slide_width)),
        height=Emu(round(h / img_h * prs.slide_height)),
    )
    try:
        pic.name = name
    except Exception:
        pass


def build_cutout_ppt(pptx: Path, boxes_path: Path, out_path: Path, work_dir: Path, overwrite: bool = False, keep_assets: bool = False) -> Path:
    payload = read_json(boxes_path)
    slides_payload = payload["slides"] if isinstance(payload, dict) else payload
    boxes_by_slide = {int(item["slide"]): item.get("boxes", []) for item in slides_payload}
    out_path = unique_output_path(out_path, overwrite)
    assets_dir = work_dir / "ppt_assets"
    if assets_dir.exists():
        import shutil

        shutil.rmtree(assets_dir)
    base_dir = assets_dir / "base"
    module_dir = assets_dir / "modules"
    base_dir.mkdir(parents=True, exist_ok=True)
    module_dir.mkdir(parents=True, exist_ok=True)

    src_prs = Presentation(str(pptx))
    out_prs = Presentation()
    out_prs.slide_width = src_prs.slide_width
    out_prs.slide_height = src_prs.slide_height
    blank = out_prs.slide_layouts[6]

    manifest = []
    for slide_no, source_slide in enumerate(src_prs.slides, 1):
        picture = full_slide_picture(source_slide)
        if picture is None:
            raise RuntimeError(f"slide {slide_no} has no picture source")
        original = Image.open(BytesIO(picture.image.blob)).convert("RGB")
        img_w, img_h = original.size
        boxes = [clamp_box(box, img_w, img_h) for box in boxes_by_slide.get(slide_no, [])]

        base = original.copy()
        draw = ImageDraw.Draw(base)
        pad = max(2, round(min(img_w, img_h) * 0.0025))
        for x, y, w, h in boxes:
            draw.rectangle([max(0, x - pad), max(0, y - pad), min(img_w, x + w + pad), min(img_h, y + h + pad)], fill=(255, 255, 255))
        base_path = base_dir / f"slide_{slide_no:02d}_base_whitefill.png"
        base.save(base_path)

        slide = out_prs.slides.add_slide(blank)
        background = slide.shapes.add_picture(str(base_path), 0, 0, width=out_prs.slide_width, height=out_prs.slide_height)
        try:
            background.name = f"WHITEFILL_BASE_S{slide_no:02d}"
        except Exception:
            pass

        modules = []
        for module_no, (x, y, w, h) in enumerate(boxes, 1):
            module_path = module_dir / f"slide_{slide_no:02d}_module_{module_no:02d}.png"
            original.crop((x, y, x + w, y + h)).save(module_path)
            add_picture(slide, module_path, x, y, w, h, img_w, img_h, out_prs, f"MODULE_S{slide_no:02d}_{module_no:02d}")
            modules.append({"index": module_no, "x": x, "y": y, "w": w, "h": h, "file": str(module_path)})
        manifest.append({"slide": slide_no, "source_size": [img_w, img_h], "modules": modules})

    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_prs.save(str(out_path))
    write_json(work_dir / "module_cutouts_build_manifest.json", {"source": str(pptx), "output": str(out_path), "slides": manifest})
    if not keep_assets:
        import shutil

        shutil.rmtree(assets_dir, ignore_errors=True)
    return out_path


def parse_args(argv: Iterable[str] | None = None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    sub = parser.add_subparsers(dest="command", required=True)

    prepare = sub.add_parser("prepare", help="extract slide images and auto-detect module boxes")
    prepare.add_argument("--pptx", required=True, type=Path)
    prepare.add_argument("--work-dir", required=True, type=Path)
    prepare.add_argument("--force", action="store_true")

    serve = sub.add_parser("serve", help="serve the red-box adjustment page")
    serve.add_argument("--work-dir", required=True, type=Path)
    serve.add_argument("--port", type=int, default=8791)

    review = sub.add_parser("review", help="prepare, then serve the red-box adjustment page")
    review.add_argument("--pptx", required=True, type=Path)
    review.add_argument("--work-dir", required=True, type=Path)
    review.add_argument("--port", type=int, default=8791)
    review.add_argument("--force", action="store_true")

    build = sub.add_parser("build", help="build PPTX with white-filled base plus module image layers")
    build.add_argument("--pptx", required=True, type=Path)
    build.add_argument("--boxes", required=True, type=Path, help="adjusted_modules.json from the review page")
    build.add_argument("--out", required=True, type=Path)
    build.add_argument("--work-dir", required=True, type=Path)
    build.add_argument("--overwrite", action="store_true")
    build.add_argument("--keep-assets", action="store_true")
    return parser.parse_args(argv)


def main(argv: Iterable[str] | None = None) -> int:
    args = parse_args(argv)
    if args.command == "prepare":
        path = prepare_review(args.pptx, args.work_dir, args.force)
        print(path)
    elif args.command == "serve":
        serve_review(args.work_dir, args.port)
    elif args.command == "review":
        prepare_review(args.pptx, args.work_dir, args.force)
        serve_review(args.work_dir, args.port)
    elif args.command == "build":
        out = build_cutout_ppt(args.pptx, args.boxes, args.out, args.work_dir, args.overwrite, args.keep_assets)
        print(out)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
